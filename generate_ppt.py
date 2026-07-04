#!/usr/bin/env python3
"""生成三星事业部运营管理平台 培训演示PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
import copy

# ========== 配色方案 ==========
SAMSUNG_BLUE = RGBColor(0x14, 0x28, 0xA0)     # 三星深蓝
SAMSUNG_LIGHT = RGBColor(0x3F, 0x63, 0xD2)     # 中蓝
SAMSUNG_SKY = RGBColor(0x61, 0x8B, 0xF2)       # 浅蓝
ACCENT_GOLD = RGBColor(0xF2, 0x9C, 0x38)       # 金色强调
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)       # 绿色
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)         # 红色
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)      # 橙色
ACCENT_PURPLE = RGBColor(0x8E, 0x44, 0xAD)      # 紫色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
BORDER_GRAY = RGBColor(0xDC, 0xDF, 0xE4)
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFC)

# ========== 辅助函数 ==========
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 页面边距
MARGIN_L = Inches(0.8)
MARGIN_R = Inches(0.8)
CONTENT_W = Inches(11.733)


def add_blank_slide():
    """添加空白幻灯片"""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_rect(slide, left, top, width, height, fill_color, border_color=None, corner_radius=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=Pt(14), color=BLACK,
                bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, default_size=Pt(14),
                          default_color=BLACK, default_bold=False, alignment=PP_ALIGN.LEFT):
    """添加多行文本框，每行一个tuple: (text, size, color, bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, size, color, bold = line_data, default_size, default_color, default_bold
        else:
            text = line_data[0]
            size = line_data[1] if len(line_data) > 1 else default_size
            color = line_data[2] if len(line_data) > 2 else default_color
            bold = line_data[3] if len(line_data) > 3 else default_bold
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Microsoft YaHei'
        p.alignment = alignment
        p.space_after = Pt(4)
    return txBox


def add_arrow(slide, start_left, start_top, end_left, end_top, color=SAMSUNG_LIGHT, width=Pt(2)):
    """添加连接线箭头"""
    connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, start_left, start_top, end_left, end_top)
    connector.line.color.rgb = color
    connector.line.width = width
    # 添加箭头
    connector.line._ln.append(
        qn('a:tailEnd')  # This sets arrowhead formatting
    )
    return connector


def add_flow_arrow(slide, left, top, width, height, fill_color=SAMSUNG_LIGHT):
    """添加流程箭头形状"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color, text='', text_color=WHITE, font_size=Pt(12)):
    """添加圆形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = text_color
        p.font.bold = True
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
    return shape


def add_slide_number(slide, num, total=21):
    """添加页码"""
    add_textbox(slide, Inches(12.0), Inches(7.05), Inches(1.2), Inches(0.35),
                f'{num} / {total}', Pt(10), MED_GRAY, alignment=PP_ALIGN.RIGHT)


def add_header_bar(slide, title_text, subtitle_text=''):
    """添加统一的页面标题栏"""
    # 顶部蓝色横条
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), SAMSUNG_BLUE)
    # 白色细线
    add_rect(slide, Inches(0), Inches(1.1), prs.slide_width, Pt(3), ACCENT_GOLD)
    # 标题
    add_textbox(slide, MARGIN_L, Inches(0.15), Inches(10), Inches(0.6),
                title_text, Pt(28), WHITE, bold=True)
    if subtitle_text:
        add_textbox(slide, MARGIN_L, Inches(0.65), Inches(10), Inches(0.4),
                    subtitle_text, Pt(14), RGBColor(0xBB, 0xCC, 0xEE))


def add_card(slide, left, top, width, height, title, content_lines,
             icon='', accent_color=SAMSUNG_LIGHT, title_color=WHITE):
    """添加卡片式内容块"""
    # 卡片背景
    card = add_rect(slide, left, top, width, height, WHITE, BORDER_GRAY, corner_radius=None)
    card.shadow.inherit = False
    # 彩色标题头
    add_rect(slide, left, top, width, Inches(0.45), accent_color)
    # 图标 + 标题 (確保 icon 是字串)
    if icon and isinstance(icon, str):
        label = icon + ' ' + title
    else:
        label = title
    add_textbox(slide, left + Inches(0.15), top + Inches(0.05), width - Inches(0.3), Inches(0.38),
                label, Pt(13), title_color, bold=True)
    # 内容
    body = '\n'.join(content_lines)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.55), width - Inches(0.3),
                height - Inches(0.65), body, Pt(11), DARK_GRAY)


def add_badge(slide, left, top, text, color=SAMSUNG_LIGHT):
    """添加标签徽章"""
    shape = add_rect(slide, left, top, Inches(1.0), Inches(0.3), color, corner_radius=None)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_table(slide, left, top, width, headers, rows, col_widths=None, header_color=SAMSUNG_BLUE):
    """添加格式化表格"""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    row_height = Inches(0.38)
    table_height = row_height * num_rows
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, table_height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # 表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = 'Microsoft YaHei'
            p.alignment = PP_ALIGN.CENTER

    # 数据行 - 交替颜色
    for r_idx, row in enumerate(rows):
        bg = BG_LIGHT if r_idx % 2 == 0 else WHITE
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK_GRAY
                p.font.name = 'Microsoft YaHei'
                p.alignment = PP_ALIGN.CENTER

    return table_shape


# ============================================================
# 开始生成 PPT
# ============================================================
TOTAL_SLIDES = 21

# ━━━━━━━━━━━━━━ Slide 1: 封面 ━━━━━━━━━━━━━━
slide = add_blank_slide()
# 背景
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, SAMSUNG_BLUE)
# 装饰几何图形
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_GOLD)
add_rect(slide, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), ACCENT_GOLD)
# 大背景圆装饰
for i, (l, t, s, c) in enumerate([
    (Inches(9.5), Inches(-1.5), Inches(6), RGBColor(0x1A, 0x30, 0xB5)),
    (Inches(-2), Inches(4), Inches(5), RGBColor(0x0E, 0x1E, 0x80)),
    (Inches(10), Inches(5), Inches(3), RGBColor(0x1A, 0x30, 0xB5)),
]):
    circle = add_circle(slide, l, t, s, c)
    circle.fill.fore_color.rgb = c
# 主标题
add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
            '三星事业部运营管理平台', Pt(52), WHITE, bold=True)
add_textbox(slide, Inches(1.5), Inches(2.9), Inches(8), Inches(0.6),
            '—— 全链路门店运营数字化解决方案', Pt(22), ACCENT_GOLD, bold=False)
# 分隔线
add_rect(slide, Inches(1.5), Inches(3.7), Inches(3), Pt(4), ACCENT_GOLD)
# 元信息
add_textbox(slide, Inches(1.5), Inches(4.1), Inches(6), Inches(0.4),
            '培训演示 · 使用指南 v2.1', Pt(16), RGBColor(0xBB, 0xCC, 0xEE))
add_textbox(slide, Inches(1.5), Inches(4.55), Inches(6), Inches(0.4),
            '2026年6月  |  简禹豪（管理员）', Pt(14), RGBColor(0x99, 0xAA, 0xCC))
# 版本号徽章
add_badge(slide, Inches(1.5), Inches(5.3), 'v2.1 试运营版', ACCENT_GOLD)


# ━━━━━━━━━━━━━━ Slide 2: 目录 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '📋 培训目录', '本次培训将覆盖以下内容')
add_slide_number(slide, 2, TOTAL_SLIDES)

toc_items = [
    ('01', '系统概述与架构', '系统定位、技术规格、角色权限', SAMSUNG_BLUE),
    ('02', '登录与账号体系', '登录流程、测试账号、权限说明', SAMSUNG_LIGHT),
    ('03', '📊 数据总览', 'Dashboard核心KPI、图表分析', SAMSUNG_SKY),
    ('04', '💰 销售进度', '月度报表、对比分析、门店排名', ACCENT_GOLD),
    ('05', '📦 库存监控', '库存展示、预警机制、补货建议', ACCENT_ORANGE),
    ('06', '💲 价格看板', '指导价与成交价对比', ACCENT_GREEN),
    ('07', '💵 店员提成', '提成计算、排行、活动政策', ACCENT_PURPLE),
    ('08', '🤖 AI 智能助手', 'DeepSeek驱动的业务分析', SAMSUNG_BLUE),
    ('09', '协作工具集', '社区/知识库/考勤/审批/任务', SAMSUNG_LIGHT),
    ('10', '⚙️ 后台管理', '用户管理、数据同步、系统配置', SAMSUNG_SKY),
    ('11', '常见问题 FAQ', '登录/数据/移动端/故障排查', ACCENT_RED),
]

for i, (num, title, desc, color) in enumerate(toc_items):
    col = i % 2
    row = i // 2
    left = Inches(0.8) + Inches(6.1) * col
    top = Inches(1.5) + Inches(0.98) * row

    # 编号圆形
    add_circle(slide, left, top, Inches(0.55), color, num, WHITE, Pt(16))
    # 标题
    add_textbox(slide, left + Inches(0.75), top + Inches(0.02), Inches(4.5), Inches(0.32),
                title, Pt(14), DARK_GRAY, bold=True)
    # 描述
    add_textbox(slide, left + Inches(0.75), top + Inches(0.32), Inches(4.5), Inches(0.25),
                desc, Pt(10), MED_GRAY)


# ━━━━━━━━━━━━━━ Slide 3: 系统概述 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '01 系统概述与架构', '全链路门店运营数字化管理平台')
add_slide_number(slide, 3, TOTAL_SLIDES)

# 概述卡片
add_card(slide, MARGIN_L, Inches(1.5), Inches(5.5), Inches(2.2),
         '📐 系统定位', [
             '• 面向三星授权体验店的运营管理平台',
             '• 覆盖 10+ 业务模块：销售/库存/价格/提成/会员/知识库/协作/AI',
             '• 支持 10 家门店同时运营',
             '• B/S 架构，浏览器即用，无需安装客户端',
         ], SAMSUNG_BLUE)

add_card(slide, Inches(6.7), Inches(1.5), Inches(5.5), Inches(2.2),
         '⚙️ 技术规格', [
             '• 访问地址：http://localhost:9527',
             '• 推荐浏览器：Chrome / Edge (推荐 Chrome)',
             '• 数据存储：SQLite 本地数据库',
             '• 会话时长：登录后 72 小时有效',
             '• 支持移动端 + PWA 安装到桌面',
         ], SAMSUNG_LIGHT)

# 角色权限表
add_textbox(slide, MARGIN_L, Inches(4.0), Inches(4), Inches(0.35),
            '👥 角色与权限矩阵', Pt(16), DARK_GRAY, bold=True)

role_headers = ['角色', '权限范围']
role_rows = [
    ['管理员 (admin)', '全部功能 + 用户管理 + 系统配置'],
    ['店长 (manager)', '本店数据 + 通用提成/知识库/协作'],
    ['店员 (staff)', '有限访问（打卡/知识库查看等）'],
]
add_table(slide, MARGIN_L, Inches(4.45), Inches(11.5), role_headers, role_rows,
          col_widths=[Inches(2.5), Inches(9.0)])

# 架构示意图 - 用方框+箭头展示
arch_y = Inches(5.9)
modules = ['浏览器\n前端', 'FastAPI\n后端', 'SQLite\n数据库', 'DeepSeek\nAI']
colors = [SAMSUNG_SKY, SAMSUNG_LIGHT, SAMSUNG_BLUE, ACCENT_GOLD]
for i, (mod, c) in enumerate(zip(modules, colors)):
    l = Inches(2.2) + Inches(2.5) * i
    shape = add_rect(slide, l, arch_y, Inches(2.0), Inches(0.85), c, corner_radius=None)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = mod
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    # 箭头
    if i < len(modules) - 1:
        arr = add_rect(slide, l + Inches(2.0), arch_y + Inches(0.35), Inches(0.5), Pt(4), SAMSUNG_LIGHT)


# ━━━━━━━━━━━━━━ Slide 4: 登录指南 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '02 登录与账号体系', '登录流程 | 测试账号 | 安全提醒')
add_slide_number(slide, 4, TOTAL_SLIDES)

# 登录三步曲
steps = [
    ('1️⃣', '打开浏览器', '访问 http://localhost:9527', SAMSUNG_BLUE),
    ('2️⃣', '输入账号密码', '输入分配的用户名和密码', SAMSUNG_LIGHT),
    ('3️⃣', '点击登录', '验证成功后进入数据总览', SAMSUNG_SKY),
]
for i, (emoji, title, desc, color) in enumerate(steps):
    left = Inches(1.2) + Inches(4.0) * i
    add_circle(slide, left + Inches(0.9), Inches(1.5), Inches(1.0), color, emoji, WHITE, Pt(28))
    add_textbox(slide, left, Inches(2.7), Inches(2.8), Inches(0.35),
                title, Pt(16), DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(3.05), Inches(2.8), Inches(0.3),
                desc, Pt(11), MED_GRAY, alignment=PP_ALIGN.CENTER)
    if i < 2:
        add_rect(slide, left + Inches(2.85), Inches(2.0), Inches(1.1), Pt(3), ACCENT_GOLD)

# 测试账号表
add_textbox(slide, MARGIN_L, Inches(3.6), Inches(4), Inches(0.35),
            '🔑 测试账号（管理员）', Pt(14), DARK_GRAY, bold=True)

admin_headers = ['账号', '密码', '角色', '说明']
admin_rows = [
    ['admin', 'admin123', '管理员', '超级管理员，可查看所有门店'],
    ['assistant', 'admin123', '管理员', '运营助理，同管理员权限'],
]
add_table(slide, MARGIN_L, Inches(3.95), Inches(5.5), admin_headers, admin_rows,
          col_widths=[Inches(1.0), Inches(1.0), Inches(1.0), Inches(2.5)])

# 店长账号（缩略）
add_textbox(slide, Inches(6.7), Inches(3.6), Inches(4), Inches(0.35),
            '🔑 测试账号（店长）', Pt(14), DARK_GRAY, bold=True)

store_headers = ['账号', '密码', '门店']
store_rows = [
    ['store1', 'store123', '万象城旗舰店'],
    ['store2', 'store123', '华润万象汇体验店'],
    ['store3', 'store123', '兴义梦乐城体验店'],
    ['store4', 'store123', '遵义吾悦体验店'],
    ['store5', 'store123', '曲靖万达体验店'],
    ['…store10', 'store123', '共10家门店'],
]
add_table(slide, Inches(6.7), Inches(3.95), Inches(5.5), store_headers, store_rows,
          col_widths=[Inches(1.2), Inches(1.0), Inches(3.3)])

# 安全提醒
warning = add_rect(slide, MARGIN_L, Inches(6.15), Inches(11.5), Inches(0.55), RGBColor(0xFF, 0xF3, 0xE0), ACCENT_ORANGE, corner_radius=None)
add_textbox(slide, Inches(1.1), Inches(6.22), Inches(10), Inches(0.4),
            '⚠️ 安全提醒：试运营结束后请立即修改默认密码！修改路径：右上角头像 → 修改密码',
            Pt(11), ACCENT_ORANGE, bold=True)

# 移动端提示
add_textbox(slide, MARGIN_L, Inches(6.7), Inches(11.5), Inches(0.35),
            '💡 移动端：手机浏览器访问 http://<电脑IP>:9527 → 浏览器菜单 → "添加到主屏幕" → 安装为 PWA 应用',
            Pt(11), MED_GRAY)


# ━━━━━━━━━━━━━━ Slide 5: 数据总览 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '03 📊 数据总览 Dashboard', '登录后默认首页 | 核心KPI一目了然')
add_slide_number(slide, 5, TOTAL_SLIDES)

# 页面结构图
dashboard_structure = [
    ('统计卡片区', '门店数 | 月销售额 | 月完成率 | 库存总量', SAMSUNG_BLUE, Inches(3.5)),
    ('图表分析区', '销售趋势折线图 | 机型占比饼图 | 门店排行柱状图', SAMSUNG_LIGHT, Inches(2.5)),
    ('快捷入口区', '一键跳转：销售进度 · 库存监控 · 店员提成等', SAMSUNG_SKY, Inches(1.5)),
]
for i, (title, desc, color, h) in enumerate(dashboard_structure):
    y = Inches(1.6) + Inches(0.3)
    if i > 0:
        y = Inches(1.6) + Inches(3.8) if i == 1 else Inches(1.6) + Inches(6.2) - Inches(0.2)
    if i == 0:
        y = Inches(1.6)
        left = Inches(1.8)
        w = Inches(9.5)
    elif i == 1:
        y = Inches(3.4)
        left = Inches(1.8)
        w = Inches(9.5)
    else:
        y = Inches(6.0)
        left = Inches(1.8)
        w = Inches(9.5)

    rect = add_rect(slide, left, y, w, Inches(1.1), color, corner_radius=None)
    add_textbox(slide, left + Inches(0.3), y + Inches(0.1), w - Inches(0.6), Inches(0.35),
                title, Pt(16), WHITE, bold=True)
    add_textbox(slide, left + Inches(0.3), y + Inches(0.55), w - Inches(0.6), Inches(0.35),
                desc, Pt(11), RGBColor(0xDD, 0xDF, 0xF5))

# KPI 卡片示意
kpi_items = [
    ('🏪', '10', '门店数', SAMSUNG_BLUE),
    ('💰', '¥XXX万', '月销售额', ACCENT_GREEN),
    ('📊', 'XX%', '月完成率', ACCENT_GOLD),
    ('📦', 'XXX台', '库存总量', SAMSUNG_SKY),
]
for i, (icon, val, label, color) in enumerate(kpi_items):
    cx = Inches(2.2) + Inches(2.9) * i
    card = add_rect(slide, cx, Inches(1.8), Inches(2.4), Inches(1.25), WHITE, BORDER_GRAY, corner_radius=None)
    add_textbox(slide, cx, Inches(1.85), Inches(2.4), Inches(0.3),
                icon, Pt(22), color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, cx, Inches(2.2), Inches(2.4), Inches(0.4),
                val, Pt(22), DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, cx, Inches(2.55), Inches(2.4), Inches(0.3),
                label, Pt(10), MED_GRAY, alignment=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━ Slide 6: 销售进度 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '04 💰 销售进度', '月度报表 | 对比分析 | 门店排名')
add_slide_number(slide, 6, TOTAL_SLIDES)

# 三个Tab卡片
tabs = [
    ('📊 月度报表', ['• 各门店手机销售额/台量', '• NCME销售额', '• 配件销售额', '• 月度完成率'], SAMSUNG_BLUE),
    ('📈 对比分析', ['• 环比（vs上月）', '• 同比（vs去年同月）', '• Chart.js柱状图', '• 销售趋势可视化'], ACCENT_GREEN),
    ('🏆 门店排名', ['• 按完成率排名', '• 按台量排名', '• 🥇🥈🥉 奖牌展示', '• 支持导出数据'], ACCENT_GOLD),
]
for i, (title, items, color) in enumerate(tabs):
    left = MARGIN_L + Inches(4.1) * i
    add_card(slide, left, Inches(1.5), Inches(3.8), Inches(2.3), title, items, accent_color=color)

# 指标说明表
add_textbox(slide, MARGIN_L, Inches(4.1), Inches(4), Inches(0.35),
            '📐 关键指标说明', Pt(14), DARK_GRAY, bold=True)

ind_headers = ['指标', '含义']
ind_rows = [
    ['手机销售额', 'S26/Fold7/Flip7/W26 等机型销售额合计'],
    ['台量', '销售手机总台数'],
    ['NCME 销售额', 'NCME 类产品销售额'],
    ['配件销售额', '手机壳/贴膜/耳机/充电器等'],
    ['完成率', '销售额 ÷ 月度目标 × 100%'],
]
add_table(slide, MARGIN_L, Inches(4.45), Inches(11.5), ind_headers, ind_rows,
          col_widths=[Inches(2.5), Inches(9.0)])

# 操作流程
add_textbox(slide, MARGIN_L, Inches(6.3), Inches(4), Inches(0.35),
            '⚡ 操作流程', Pt(14), DARK_GRAY, bold=True)

flow_steps = ['选择月份', '切换门店', '查看报表', '导出数据']
for i, step in enumerate(flow_steps):
    l = Inches(1.5) + Inches(2.8) * i
    add_rect(slide, l, Inches(6.75), Inches(2.2), Inches(0.5), SAMSUNG_LIGHT, corner_radius=None)
    add_textbox(slide, l, Inches(6.78), Inches(2.2), Inches(0.4),
                step, Pt(11), WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(flow_steps) - 1:
        arr = add_rect(slide, l + Inches(2.2), Inches(6.93), Inches(0.6), Pt(3), ACCENT_GOLD)


# ━━━━━━━━━━━━━━ Slide 7: 库存监控 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '05 📦 库存监控', '库存展示 | 预警机制 | 补货建议')
add_slide_number(slide, 7, TOTAL_SLIDES)

# 库存状态三色灯
status_items = [
    ('🟢 正常', '库存充足', '可销 ≥ 7天', '维持现状', ACCENT_GREEN),
    ('🟡 预警', '库存紧张', '可销 < 7天', '关注并准备补货', ACCENT_ORANGE),
    ('🔴 紧急', '库存告急', '可销 < 3天', '立即下单补货', ACCENT_RED),
]
for i, (label, status, condition, action, color) in enumerate(status_items):
    left = Inches(1.0) + Inches(4.0) * i
    card = add_rect(slide, left, Inches(1.6), Inches(3.5), Inches(1.8), WHITE, BORDER_GRAY, corner_radius=None)
    # 状态灯
    add_circle(slide, left + Inches(1.25), Inches(1.75), Inches(0.6), color, label[:2], WHITE, Pt(12))
    add_textbox(slide, left, Inches(2.5), Inches(3.5), Inches(0.35),
                label, Pt(14), DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(2.8), Inches(3.5), Inches(0.25),
                f'{condition}', Pt(10), MED_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(3.05), Inches(3.5), Inches(0.25),
                f'建议：{action}', Pt(10), color, bold=True, alignment=PP_ALIGN.CENTER)

# 功能模块
add_textbox(slide, MARGIN_L, Inches(3.7), Inches(4), Inches(0.35),
            '📋 库存管理功能', Pt(14), DARK_GRAY, bold=True)

inv_headers = ['功能模块', '说明']
inv_rows = [
    ['库存列表', '查看各门店/各机型当前库存数量和状态'],
    ['库存预警', '可销天数不足时自动标红/标黄提醒'],
    ['补货建议', '智能分析销售速度，给出补货优先级建议'],
    ['门店筛选', '按门店筛选查看库存分布'],
    ['机型筛选', '按S26/Fold7/Flip7等机型系列筛选'],
    ['数据同步', '支持从亿博士（eBoss）系统自动同步'],
]
add_table(slide, MARGIN_L, Inches(4.05), Inches(11.5), inv_headers, inv_rows,
          col_widths=[Inches(2.5), Inches(9.0)])


# ━━━━━━━━━━━━━━ Slide 8: 价格看板 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '06 💲 价格看板', '指导价与成交价对比分析')
add_slide_number(slide, 8, TOTAL_SLIDES)

add_card(slide, MARGIN_L, Inches(1.5), Inches(5.5), Inches(2.0),
         '📊 价格看板功能', [
             '• 展示各机型官方指导价与实际成交价对比',
             '• 价差列：红色负值 / 绿色正值',
             '• 帮助店长掌握价格政策执行情况',
             '• 及时发现异常定价行为',
         ], SAMSUNG_BLUE)

add_card(slide, Inches(6.7), Inches(1.5), Inches(5.5), Inches(2.0),
         '📐 数据列说明', [
             '• 机型：S26 / Fold7 / Flip7 / W26 等',
             '• 官方指导价：三星官方建议零售价',
             '• 实际成交均价：门店实际销售均价',
             '• 价差：实际价 − 指导价',
         ], ACCENT_GREEN)

# 示例数据表
add_textbox(slide, MARGIN_L, Inches(3.8), Inches(4), Inches(0.35),
            '📋 数据示例', Pt(14), DARK_GRAY, bold=True)

price_headers = ['机型', '官方指导价', '实际成交均价', '价差', '状态']
price_rows = [
    ['Galaxy S26', '¥6,999', '¥6,800', '-¥199', '🟢 正常'],
    ['Galaxy Fold7', '¥12,999', '¥12,500', '-¥499', '🟢 正常'],
    ['Galaxy Flip7', '¥8,999', '¥8,600', '-¥399', '🟢 正常'],
    ['Galaxy W26', '¥10,999', '¥10,200', '-¥799', '🟡 关注'],
    ['Galaxy A56', '¥3,299', '¥3,100', '-¥199', '🟢 正常'],
]
add_table(slide, MARGIN_L, Inches(4.15), Inches(11.5), price_headers, price_rows,
          col_widths=[Inches(2.2), Inches(2.2), Inches(2.2), Inches(1.9), Inches(3.0)])


# ━━━━━━━━━━━━━━ Slide 9: 店员提成 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '07 💵 店员提成', '提成计算 | 排行 | 规则 | 活动政策')
add_slide_number(slide, 9, TOTAL_SLIDES)

# 5个Tab
tabs2 = [
    ('👥 店员列表', '按门店查看店员底薪/销量/提成/收入'),
    ('📊 提成概览', 'Chart.js环形图展示提成分布'),
    ('🏆 销售排行', '🥇🥈🥉 销冠榜单'),
    ('📋 提成规则', '各品类提成标准定价'),
    ('📢 活动政策', '当前活动政策和培训通知'),
]
for i, (title, desc) in enumerate(tabs2):
    left = MARGIN_L + Inches(6.1) * (i % 2)
    top = Inches(1.5) + Inches(1.0) * (i // 2)
    shape = add_rect(slide, left, top, Inches(5.8), Inches(0.75), SAMSUNG_LIGHT, corner_radius=None)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.08), Inches(5.4), Inches(0.3),
                title, Pt(13), WHITE, bold=True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.38), Inches(5.4), Inches(0.3),
                desc, Pt(10), RGBColor(0xDD, 0xDF, 0xF5))

# 提成公式
add_textbox(slide, MARGIN_L, Inches(5.1), Inches(8), Inches(0.35),
            '🧮 提成计算公式', Pt(14), DARK_GRAY, bold=True)

formula = add_rect(slide, MARGIN_L, Inches(5.55), Inches(11.5), Inches(1.0),
                   RGBColor(0xEE, 0xF1, 0xFA), SAMSUNG_LIGHT, corner_radius=None)
add_textbox(slide, Inches(1.2), Inches(5.65), Inches(10.5), Inches(0.8),
            '预计收入  =  底薪  +  手机销量×手机提成单价  +  NCME销售额×NCME提成比例\n'
            '                    +  配件销售额×配件提成比例  +  回收提成',
            Pt(13), SAMSUNG_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━ Slide 10: 会员管理 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '08 👥 会员管理', '会员信息管理 · 搜索筛选 · 新增会员')
add_slide_number(slide, 10, TOTAL_SLIDES)

# 功能卡片
mem_functions = [
    ('👥 会员列表', ['查看会员信息', '姓名/手机号/等级', '累计消费/注册门店', '支持搜索筛选'], SAMSUNG_BLUE),
    ('🔍 搜索筛选', ['按姓名搜索', '按手机号搜索', '按门店筛选', '实时过滤结果'], SAMSUNG_LIGHT),
    ('➕ 新增会员', ['手动录入新会员', '填写基本信息', '选择门店归属', '即时生效'], ACCENT_GREEN),
]
for i, (title, items, color) in enumerate(mem_functions):
    left = MARGIN_L + Inches(4.1) * i
    add_card(slide, left, Inches(1.5), Inches(3.8), Inches(2.3), title, items, accent_color=color)

# 会员等级
add_textbox(slide, MARGIN_L, Inches(4.1), Inches(4), Inches(0.35),
            '⭐ 会员等级体系', Pt(14), DARK_GRAY, bold=True)

level_items = [
    ('🥉 普通会员', '注册即有', SAMSUNG_SKY),
    ('🥈 银卡会员', '累计消费 ≥ ¥5,000', SAMSUNG_LIGHT),
    ('🥇 金卡会员', '累计消费 ≥ ¥20,000', ACCENT_GOLD),
    ('💎 钻石会员', '累计消费 ≥ ¥50,000', SAMSUNG_BLUE),
]
for i, (name, cond, color) in enumerate(level_items):
    left = Inches(1.5) + Inches(2.8) * i
    rect = add_rect(slide, left, Inches(4.55), Inches(2.4), Inches(1.0), color, corner_radius=None)
    add_textbox(slide, left, Inches(4.6), Inches(2.4), Inches(0.35),
                name, Pt(14), WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(4.95), Inches(2.4), Inches(0.3),
                cond, Pt(9), RGBColor(0xEE, 0xEE, 0xFF), alignment=PP_ALIGN.CENTER)

# 权限说明
add_rect(slide, MARGIN_L, Inches(5.9), Inches(11.5), Inches(0.55), RGBColor(0xFF, 0xF3, 0xE0), ACCENT_ORANGE, corner_radius=None)
add_textbox(slide, Inches(1.1), Inches(5.97), Inches(11), Inches(0.4),
            '🔒 权限说明：仅管理员和店长角色可访问会员管理模块。管理员可跨店查看，店长仅看本店会员。',
            Pt(11), ACCENT_ORANGE, bold=True)


# ━━━━━━━━━━━━━━ Slide 11: 店长百事通 + 交流社区 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '09 📚 店长百事通 & 💬 交流社区', '知识库 + 团队协作双引擎')
add_slide_number(slide, 11, TOTAL_SLIDES)

# 店长百事通
add_card(slide, MARGIN_L, Inches(1.5), Inches(5.5), Inches(2.5),
         '📚 店长百事通 - 知识库', [
             '• 分类浏览：产品知识/销售技巧/政策法规/培训资料',
             '• 全文搜索：输入关键词搜索所有文章内容',
             '• 文章详情：查看全文 + 自动生成摘要',
             '• 使用场景：查参数/学话术/看政策/查培训',
         ], SAMSUNG_BLUE)

# 交流社区
add_card(slide, Inches(6.7), Inches(1.5), Inches(5.5), Inches(2.5),
         '💬 交流社区 - 团队协作', [
             '• 发帖：发布话题（文字+图片）+选择分类',
             '• 回复：在帖子下评论交流互动',
             '• 分类：经验分享/问题求助/公告通知',
             '• 典型场景：分享经验、求助问题、发布通知',
         ], ACCENT_GREEN)

# 知识库分类标签
add_textbox(slide, MARGIN_L, Inches(4.3), Inches(4), Inches(0.35),
            '🏷️ 知识库分类', Pt(14), DARK_GRAY, bold=True)

cats = ['产品知识', '销售技巧', '政策法规', '培训资料', '竞品分析', '常见问题']
for i, cat in enumerate(cats):
    left = Inches(1.0) + Inches(2.0) * i
    add_badge(slide, left, Inches(4.75), cat, SAMSUNG_LIGHT if i % 2 == 0 else SAMSUNG_SKY)

# 使用场景示例
add_textbox(slide, MARGIN_L, Inches(5.3), Inches(8), Inches(0.35),
            '💡 典型使用场景', Pt(14), DARK_GRAY, bold=True)

scenario_data = [
    ('🔍 查参数', 'S26摄像头像素？续航多久？'),
    ('💬 学话术', '如何介绍Fold7折叠屏优势？'),
    ('📋 看政策', '最新以旧换新政策是什么？'),
    ('📖 查培训', '新品培训PPT在哪里？'),
    ('📢 分享经验', '"今天W26卖了8台，分享下心得"'),
    ('❓ 求助问题', '"Flip7外屏有客户反馈花屏？"'),
]
for i, (action, example) in enumerate(scenario_data):
    col = i % 3
    row = i // 3
    left = Inches(1.0) + Inches(4.0) * col
    top = Inches(5.75) + Inches(0.5) * row
    add_textbox(slide, left, top, Inches(3.8), Inches(0.35),
                f'{action}  →  {example}', Pt(10), DARK_GRAY)


# ━━━━━━━━━━━━━━ Slide 12: AI 智能助手 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '10 🤖 AI 智能助手', 'DeepSeek 驱动的智能业务分析')
add_slide_number(slide, 12, TOTAL_SLIDES)

# AI能力网格
ai_abilities = [
    ('📊 销售查询', '万象城店本月卖了多少台S26？', SAMSUNG_BLUE),
    ('📦 库存查询', '哪个门店W26库存最少？需要补货吗？', SAMSUNG_LIGHT),
    ('📚 知识检索', '以旧换新政策是什么？', SAMSUNG_SKY),
    ('📈 数据分析', '对比万象城和万象汇的本月完成率？', ACCENT_GOLD),
    ('💵 提成查询', '本月店员提成最高的是谁？', ACCENT_GREEN),
    ('📝 综合建议', '帮我分析哪些门店需要重点关注？', ACCENT_PURPLE),
]
for i, (ability, example, color) in enumerate(ai_abilities):
    col = i % 2
    row = i // 2
    left = Inches(1.0) + Inches(6.1) * col
    top = Inches(1.5) + Inches(1.3) * row
    card = add_rect(slide, left, top, Inches(5.8), Inches(1.05), WHITE, BORDER_GRAY, corner_radius=None)
    # 左侧色条
    add_rect(slide, left, top, Inches(0.08), Inches(1.05), color)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.08), Inches(5.0), Inches(0.35),
                ability, Pt(14), DARK_GRAY, bold=True)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.55), Inches(5.0), Inches(0.35),
                f'"…{example}…"', Pt(11), MED_GRAY)

# AI使用方式
add_textbox(slide, MARGIN_L, Inches(5.7), Inches(4), Inches(0.35),
            '⚡ AI 助手使用流程', Pt(14), DARK_GRAY, bold=True)

ai_flow = ['输入问题', 'AI 流式\n逐字输出', '自动调用\n后端工具', '返回分析\n结果']
for i, step in enumerate(ai_flow):
    l = Inches(1.5) + Inches(2.8) * i
    shape = add_rect(slide, l, Inches(6.15), Inches(2.2), Inches(0.8), SAMSUNG_LIGHT, corner_radius=None)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    if i < len(ai_flow) - 1:
        add_rect(slide, l + Inches(2.2), Inches(6.5), Inches(0.6), Pt(3), ACCENT_GOLD)


# ━━━━━━━━━━━━━━ Slide 13: 考勤打卡 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '11 📅 考勤打卡', '上班打卡 | 下班打卡 | 月度汇总')
add_slide_number(slide, 13, TOTAL_SLIDES)

# 打卡界面元素
# 实时时钟
clock = add_rect(slide, Inches(4.0), Inches(1.5), Inches(5.3), Inches(1.2), SAMSUNG_BLUE, corner_radius=None)
add_textbox(slide, Inches(4.0), Inches(1.55), Inches(5.3), Inches(0.4),
            '🕐 实时时钟', Pt(14), RGBColor(0xBB, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(4.0), Inches(2.0), Inches(5.3), Inches(0.55),
            '2026年6月8日  星期一   15:11:05', Pt(22), WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 打卡按钮示意
btn_up = add_rect(slide, Inches(2.5), Inches(3.1), Inches(3.5), Inches(1.0), ACCENT_GREEN, corner_radius=None)
add_textbox(slide, Inches(2.5), Inches(3.2), Inches(3.5), Inches(0.8),
            '✅  上班打卡\n09:00', Pt(18), WHITE, bold=True, alignment=PP_ALIGN.CENTER)

btn_down = add_rect(slide, Inches(7.3), Inches(3.1), Inches(3.5), Inches(1.0), ACCENT_ORANGE, corner_radius=None)
add_textbox(slide, Inches(7.3), Inches(3.2), Inches(3.5), Inches(0.8),
            '🔚  下班打卡\n18:00', Pt(18), WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 考勤统计
add_textbox(slide, MARGIN_L, Inches(4.4), Inches(4), Inches(0.35),
            '📊 月度考勤汇总', Pt(14), DARK_GRAY, bold=True)

att_headers = ['统计项', '说明']
att_rows = [
    ['出勤天数', '当月正常出勤总天数'],
    ['迟到次数', '超过上班时间打卡的次数'],
    ['早退次数', '早于下班时间打卡的次数'],
    ['缺卡次数', '未打卡的工作日天数'],
    ['加班时长', '超出正常工作时长的累计时间'],
]
add_table(slide, MARGIN_L, Inches(4.75), Inches(7.0), att_headers, att_rows,
          col_widths=[Inches(2.0), Inches(5.0)])


# ━━━━━━━━━━━━━━ Slide 14: 审批管理 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '12 📝 审批管理', '在线提交和处理各类申请')
add_slide_number(slide, 14, TOTAL_SLIDES)

# 审批类型
approval_types = [
    ('🏖️ 请假', '事假/病假/年假', SAMSUNG_SKY),
    ('💰 报销', '差旅/招待/办公用品', ACCENT_GREEN),
    ('🔄 调拨', '门店间库存调拨', ACCENT_ORANGE),
    ('📋 其他', '自定义审批事项', SAMSUNG_LIGHT),
]
for i, (title, desc, color) in enumerate(approval_types):
    left = Inches(1.0) + Inches(3.0) * i
    card = add_rect(slide, left, Inches(1.5), Inches(2.6), Inches(1.3), WHITE, BORDER_GRAY, corner_radius=None)
    add_rect(slide, left, Inches(1.5), Inches(2.6), Inches(0.5), color)
    add_textbox(slide, left, Inches(1.55), Inches(2.6), Inches(0.35),
                title, Pt(16), WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(2.1), Inches(2.6), Inches(0.3),
                desc, Pt(10), MED_GRAY, alignment=PP_ALIGN.CENTER)

# 审批流程
add_textbox(slide, MARGIN_L, Inches(3.1), Inches(4), Inches(0.35),
            '⚡ 审批流程', Pt(14), DARK_GRAY, bold=True)

flow_items = ['新建申请', '选择类型\n填写内容', '选择审批人\n默认上级', '提交等待\n审批处理', '通过/驳回\n通知申请人']
for i, step in enumerate(flow_items):
    l = Inches(0.6) + Inches(2.5) * i
    shape = add_rect(slide, l, Inches(3.55), Inches(2.0), Inches(0.85), SAMSUNG_LIGHT, corner_radius=None)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    if i < 4:
        add_rect(slide, l + Inches(2.0), Inches(3.93), Inches(0.5), Pt(3), ACCENT_GOLD)

# 注意事项
add_textbox(slide, MARGIN_L, Inches(4.8), Inches(11.5), Inches(0.35),
            '💡 操作提示', Pt(14), DARK_GRAY, bold=True)
tips = [
    '• 提交申请：点击「新建申请」选择类型 → 填写标题和详情 → 选择审批人 → 提交',
    '• 审批处理：切换到「待审批」Tab → 查看申请详情 → 点击「通过」或「驳回」',
    '• 状态追踪：可在「我的申请」中查看所有已提交的审批进度',
]
for i, tip in enumerate(tips):
    add_textbox(slide, Inches(1.2), Inches(5.2) + Inches(0.35) * i, Inches(10), Inches(0.3),
                tip, Pt(11), DARK_GRAY)


# ━━━━━━━━━━━━━━ Slide 15: 任务管理 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '13 ✅ 任务管理', '团队任务分配与状态追踪')
add_slide_number(slide, 15, TOTAL_SLIDES)

# 任务状态流转
add_textbox(slide, MARGIN_L, Inches(1.5), Inches(4), Inches(0.35),
            '🔄 任务状态流转', Pt(16), DARK_GRAY, bold=True)

states = [
    ('待处理', SAMSUNG_SKY),
    ('进行中', ACCENT_GOLD),
    ('已完成', ACCENT_GREEN),
]
for i, (state, color) in enumerate(states):
    l = Inches(2.2) + Inches(3.5) * i
    shape = add_rect(slide, l, Inches(2.1), Inches(2.5), Inches(1.0), color, corner_radius=None)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = state
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    if i < 2:
        add_rect(slide, l + Inches(2.5), Inches(2.55), Inches(1.0), Pt(4), ACCENT_GOLD)

# 优先级
add_textbox(slide, MARGIN_L, Inches(3.4), Inches(4), Inches(0.35),
            '🔴🟡🟢 任务优先级', Pt(14), DARK_GRAY, bold=True)

prio_items = [
    ('🔴 紧急', '需24小时内处理', ACCENT_RED),
    ('🟡 重要', '本周内完成', ACCENT_ORANGE),
    ('🟢 普通', '按计划推进', ACCENT_GREEN),
]
for i, (label, desc, color) in enumerate(prio_items):
    left = Inches(1.5) + Inches(3.8) * i
    add_rect(slide, left, Inches(3.85), Inches(3.3), Inches(0.85), color, corner_radius=None)
    add_textbox(slide, left, Inches(3.9), Inches(3.3), Inches(0.4),
                f'{label}  |  {desc}', Pt(14), WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 操作指南
add_textbox(slide, MARGIN_L, Inches(5.0), Inches(4), Inches(0.35),
            '⚡ 操作流程', Pt(14), DARK_GRAY, bold=True)

task_ops = [
    '1. 点击「新建任务」创建任务',
    '2. 填写标题、描述、指派人、优先级、截止日期',
    '3. 被指派人可在任务列表看到新任务',
    '4. 执行中点击「开始处理」，完成后点击「标记完成」',
]
for i, op in enumerate(task_ops):
    add_textbox(slide, Inches(1.2), Inches(5.45) + Inches(0.4) * i, Inches(10), Inches(0.35),
                op, Pt(12), DARK_GRAY)


# ━━━━━━━━━━━━━━ Slide 16: 后台管理 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '14 ⚙️ 后台管理', '管理员专属 · 系统配置中心')
add_slide_number(slide, 16, TOTAL_SLIDES)

admin_modules = [
    ('👤 用户管理', '创建/禁用/删除\n用户账号', SAMSUNG_BLUE),
    ('🏪 门店管理', '添加/编辑\n门店信息', SAMSUNG_LIGHT),
    ('🔄 数据同步', '从亿博士(eBoss)\n同步销售库存', ACCENT_GREEN),
    ('⚙️ 系统配置', '修改提成规则\n预警阈值参数', ACCENT_GOLD),
    ('💾 数据库备份', '手动/自动\n备份数据库', SAMSUNG_SKY),
]
for i, (title, desc, color) in enumerate(admin_modules):
    left = Inches(0.5) + Inches(2.55) * i
    card = add_rect(slide, left, Inches(1.6), Inches(2.2), Inches(2.5), WHITE, BORDER_GRAY, corner_radius=None)
    add_rect(slide, left, Inches(1.6), Inches(2.2), Inches(0.55), color)
    add_textbox(slide, left, Inches(1.65), Inches(2.2), Inches(0.40),
                title, Pt(14), WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(2.35), Inches(2.2), Inches(0.95),
                desc, Pt(10), DARK_GRAY, alignment=PP_ALIGN.CENTER)

# 权限提示
add_rect(slide, MARGIN_L, Inches(4.5), Inches(11.5), Inches(0.55), RGBColor(0xFF, 0xF3, 0xE0), ACCENT_ORANGE, corner_radius=None)
add_textbox(slide, Inches(1.1), Inches(4.57), Inches(11), Inches(0.4),
            '🔒 后台管理仅限管理员角色访问。系统启动时预置 admin/assistant 两个管理员账号。',
            Pt(12), ACCENT_ORANGE, bold=True)

# 数据同步说明
add_textbox(slide, MARGIN_L, Inches(5.3), Inches(5), Inches(0.35),
            '🔄 数据同步机制', Pt(14), DARK_GRAY, bold=True)

sync_info = [
    ('数据来源', '亿博士 (eBoss) 系统'),
    ('同步内容', '零售单数据 + 库存汇总数据'),
    ('同步频率', '默认每小时自动同步一次'),
    ('手动触发', '管理员在后台 → 数据同步 → 手动触发'),
]
for i, (label, val) in enumerate(sync_info):
    add_textbox(slide, Inches(1.2), Inches(5.75) + Inches(0.38) * i, Inches(2.0), Inches(0.3),
                label + '：', Pt(12), DARK_GRAY, bold=True)
    add_textbox(slide, Inches(3.0), Inches(5.75) + Inches(0.38) * i, Inches(8), Inches(0.3),
                val, Pt(12), MED_GRAY)


# ━━━━━━━━━━━━━━ Slide 17-19: FAQ第1页 ━━━━━━━━━━━━━━
faq_data = [
    ('Q1：登录失败怎么办？', [
        '• 检查账号密码是否正确（区分大小写）',
        '• 确认服务是否正常运行（浏览器访问 http://localhost:9527）',
        '• 如忘记密码，联系管理员重置',
    ]),
    ('Q2：页面显示异常/排版错乱？', [
        '• 最快解决方式：按 Ctrl + Shift + R（硬刷新）跳过浏览器缓存',
        '• 如果仍有问题，清除浏览器缓存后重试',
    ]),
    ('Q3：数据不更新？', [
        '• 销售和库存数据来源于亿博士（eBoss）系统',
        '• 管理员可在后台管理 → 数据同步 → 手动触发同步',
        '• 同步频率默认每小时一次',
    ]),
    ('Q4：AI 助手不回答？', [
        '• 确认 AI API 密钥已配置（管理员操作）',
        '• 刷新页面后重试',
        '• 如果回答中包含技术标记（<| DSML |），硬刷新页面即可',
    ]),
    ('Q5：如何在手机上使用？', [
        '• 手机浏览器访问 http://<电脑IP>:9527（需同一网络）',
        '• 浏览器菜单 → "添加到主屏幕" → 安装为 PWA 应用',
        '• 安装后可在桌面像 App 一样打开',
    ]),
    ('Q6：服务突然不可用？', [
        '• 重新启动：cd D:\\Workbuudy\\samsung-ops  →  python main.py',
        '• 检查Python进程是否意外关闭',
        '• 确认 9527 端口未被占用',
        '• 检查 data/samsung_ops.db 文件是否存在且完好',
    ]),
    ('Q7：如何修改密码？', [
        '• 登录后点击右上角头像区域',
        '• 选择"修改密码"',
        '• 输入旧密码和新密码',
    ]),
]

# FAQ - 分页
for page_idx in range(2):
    slide = add_blank_slide()
    add_header_bar(slide, f'15 常见问题 FAQ', '常见问题与解决方案')
    slide_num = 17 + page_idx
    add_slide_number(slide, slide_num, TOTAL_SLIDES)

    start = page_idx * 4
    end = min(start + 4, len(faq_data))
    page_faqs = faq_data[start:end]

    for i, (question, answers) in enumerate(page_faqs):
        top = Inches(1.4) + Inches(2.85) * (i // 2)
        col = i % 2
        left = Inches(0.8) + Inches(6.1) * col
        width = Inches(5.8)

        # Q 标题
        q_bg = add_rect(slide, left, top, width, Inches(0.45), SAMSUNG_BLUE, corner_radius=None)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.05), width - Inches(0.3), Inches(0.35),
                    question, Pt(14), WHITE, bold=True)
        # A 内容
        a_bg = add_rect(slide, left, top + Inches(0.45), width,
                        Inches(0.35) * len(answers) + Inches(0.3), WHITE, BORDER_GRAY, corner_radius=None)
        for j, ans in enumerate(answers):
            add_textbox(slide, left + Inches(0.2), top + Inches(0.55) + Inches(0.32) * j,
                        width - Inches(0.4), Inches(0.28), ans, Pt(9.5), DARK_GRAY)


# ━━━━━━━━━━━━━━ Slide 20: 版本记录 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_header_bar(slide, '📋 版本记录', '系统迭代历程')
add_slide_number(slide, 20, TOTAL_SLIDES)

ver_headers = ['版本', '日期', '更新内容']
ver_rows = [
    ['v1.0', '2026-05', '初始版本：数据总览/销售/库存/价格/知识库/社区'],
    ['v2.0', '2026-06', '新增：店员提成/会员管理/AI助手/考勤/审批/任务'],
    ['v2.1', '2026-06-08', '优化：移动端适配/PWA v3/DSML过滤/排版统一/使用手册'],
]
add_table(slide, MARGIN_L, Inches(1.8), Inches(11.5), ver_headers, ver_rows,
          col_widths=[Inches(1.5), Inches(2.0), Inches(8.0)])

# 技术亮点
add_textbox(slide, MARGIN_L, Inches(3.5), Inches(5), Inches(0.35),
            '✨ 系统技术亮点', Pt(16), DARK_GRAY, bold=True)

highlights = [
    ('🎨', 'PWA 渐进式Web应用', '支持离线访问和桌面安装'),
    ('🤖', 'DeepSeek AI 集成', '自然语言查询业务数据'),
    ('📊', 'Chart.js 可视化', '丰富的图表和数据分析'),
    ('📱', '响应式设计', 'PC/平板/手机全端适配'),
    ('🔐', '角色权限控制', '管理员/店长/店员三级权限'),
    ('🔄', '亿博士数据同步', '自动采集零售单和库存数据'),
]
for i, (icon, title, desc) in enumerate(highlights):
    col = i % 3
    row = i // 3
    left = Inches(1.0) + Inches(4.0) * col
    top = Inches(4.0) + Inches(1.4) * row
    card = add_rect(slide, left, top, Inches(3.5), Inches(1.1), BG_LIGHT, BORDER_GRAY, corner_radius=None)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.08), Inches(0.5), Inches(0.4),
                icon, Pt(24), alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.6), top + Inches(0.1), Inches(2.7), Inches(0.3),
                title, Pt(13), DARK_GRAY, bold=True)
    add_textbox(slide, left + Inches(0.6), top + Inches(0.45), Inches(2.7), Inches(0.4),
                desc, Pt(10), MED_GRAY)


# ━━━━━━━━━━━━━━ Slide 21: 结束页 ━━━━━━━━━━━━━━
slide = add_blank_slide()
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, SAMSUNG_BLUE)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_GOLD)
add_rect(slide, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), ACCENT_GOLD)

# 装饰圆
for l, t, s in [(Inches(-2), Inches(4), Inches(5)), (Inches(10), Inches(-1.5), Inches(4))]:
    circle = add_circle(slide, l, t, s, RGBColor(0x1A, 0x30, 0xB5))

add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
            '感谢聆听', Pt(52), WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.6),
            '三星事业部运营管理平台', Pt(24), ACCENT_GOLD, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(3.8), Inches(2.3), Pt(4), ACCENT_GOLD)
add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.4),
            '如有疑问，请联系管理员', Pt(16), RGBColor(0xBB, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(4.7), Inches(10), Inches(0.4),
            '简禹豪 | http://localhost:9527', Pt(14), RGBColor(0x99, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)

# ============================================================
# 保存
# ============================================================
output_path = 'D:/Workbuudy/samsung-ops/三星事业部运营管理平台_培训演示_v2.1.pptx'
prs.save(output_path)
print(f'✅ PPT 已生成：{output_path}')
print(f'📊 共 {len(prs.slides)} 页幻灯片')
